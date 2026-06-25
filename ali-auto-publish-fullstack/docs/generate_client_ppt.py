# -*- coding: utf-8 -*-
"""生成「客户端」PPT（给客户看）。运行: py -3.11 docs/generate_client_ppt.py"""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parent.parent
OUTPUT_CN = ROOT / "docs" / "客户版-国际站智能运营系统介绍.pptx"
OUTPUT_EN = ROOT / "docs" / "client-deck.pptx"

PRIMARY = RGBColor(0x0F, 0x4C, 0x81)
ACCENT = RGBColor(0xFF, 0x6A, 0x00)
LIGHT_BG = RGBColor(0xF4, 0xF7, 0xFB)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK = RGBColor(0x22, 0x2A, 0x35)
MUTED = RGBColor(0x6B, 0x7A, 0x8C)
SOFT = RGBColor(0x94, 0xA3, 0xB8)


def set_slide_bg(slide, rgb: RGBColor):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = rgb


def add_title_bar(slide, title: str, subtitle: str = ""):
    bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(1.0))
    bar.fill.solid()
    bar.fill.fore_color.rgb = PRIMARY
    bar.line.fill.background()
    tf = bar.text_frame
    tf.text = title
    p = tf.paragraphs[0]
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = WHITE
    if subtitle:
        box = slide.shapes.add_textbox(Inches(0.65), Inches(1.08), Inches(8.7), Inches(0.42))
        tfb = box.text_frame
        tfb.text = subtitle
        tfb.paragraphs[0].font.size = Pt(13)
        tfb.paragraphs[0].font.color.rgb = MUTED


def add_bullets(slide, items, left=0.75, top=1.65, width=8.5, font_size=17):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(5.3))
    tf = box.text_frame
    tf.word_wrap = True
    for i, text in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = text
        p.font.size = Pt(font_size)
        p.font.color.rgb = DARK
        p.space_after = Pt(9)
        if text.startswith("  "):
            p.level = 1
            p.font.size = Pt(font_size - 2)
            p.font.color.rgb = MUTED


def add_quote_bar(slide, text: str, top=5.95):
    shape = slide.shapes.add_shape(1, Inches(0.65), Inches(top), Inches(8.7), Inches(0.72))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0xE8, 0xF2, 0xFA)
    shape.line.fill.background()
    tf = shape.text_frame
    tf.margin_left = Inches(0.2)
    tf.text = text
    p = tf.paragraphs[0]
    p.font.size = Pt(14)
    p.font.color.rgb = PRIMARY
    p.font.italic = True


def add_feature_cards(slide, cards: list[tuple[str, str]], top=1.75):
    positions = [(0.55, top), (5.15, top), (0.55, top + 2.05), (5.15, top + 2.05)]
    for (title, desc), (x, y) in zip(cards, positions[: len(cards)]):
        card = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(4.35), Inches(1.85))
        card.fill.solid()
        card.fill.fore_color.rgb = LIGHT_BG
        card.line.color.rgb = RGBColor(0xD8, 0xE2, 0xEC)
        tf = card.text_frame
        tf.margin_left = Inches(0.18)
        tf.margin_top = Inches(0.14)
        p0 = tf.paragraphs[0]
        p0.text = title
        p0.font.size = Pt(18)
        p0.font.bold = True
        p0.font.color.rgb = PRIMARY
        p1 = tf.add_paragraph()
        p1.text = desc
        p1.font.size = Pt(12)
        p1.font.color.rgb = DARK
        p1.space_before = Pt(8)


def slide_cover(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, PRIMARY)

    accent_line = slide.shapes.add_shape(1, Inches(0.75), Inches(1.85), Inches(0.12), Inches(3.8))
    accent_line.fill.solid()
    accent_line.fill.fore_color.rgb = ACCENT
    accent_line.line.fill.background()

    t1 = slide.shapes.add_textbox(Inches(1.1), Inches(2.0), Inches(8.0), Inches(1.5))
    tf = t1.text_frame
    tf.text = "国际站智能运营系统"
    tf.paragraphs[0].font.size = Pt(42)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = WHITE

    t2 = slide.shapes.add_textbox(Inches(1.1), Inches(3.35), Inches(8.0), Inches(0.9))
    tf2 = t2.text_frame
    tf2.text = "自动发品 · 数据洞察 · 一站运营"
    tf2.paragraphs[0].font.size = Pt(22)
    tf2.paragraphs[0].font.color.rgb = RGBColor(0xD0, 0xE4, 0xF4)

    t3 = slide.shapes.add_textbox(Inches(1.1), Inches(4.5), Inches(8.0), Inches(1.0))
    tf3 = t3.text_frame
    tf3.text = "专为阿里巴巴国际站卖家打造\n帮助您的团队更高效地上新、更清晰地做运营决策"
    for p in tf3.paragraphs:
        p.font.size = Pt(15)
        p.font.color.rgb = RGBColor(0xA8, 0xC0, 0xD8)

    foot = slide.shapes.add_textbox(Inches(1.1), Inches(6.2), Inches(8.0), Inches(0.4))
    foot.text_frame.text = "产品方案介绍"
    foot.text_frame.paragraphs[0].font.size = Pt(12)
    foot.text_frame.paragraphs[0].font.color.rgb = RGBColor(0x88, 0xA8, 0xC0)


def slide_about(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_title_bar(slide, "这是什么？", "一款面向国际站卖家的一体化运营工具")
    add_bullets(
        slide,
        [
            "将「批量发品、图片处理、数据下载、运营分析」整合在同一套系统中",
            "您通过电脑本地运行，在可视化界面完成配置与任务管理",
            "减少重复登录后台、手工填表、手工导数据的时间",
            "让团队把精力集中在选品、回复询盘和成交上",
        ],
        top=1.75,
        font_size=18,
    )
    add_quote_bar(slide, "一句话：帮您把国际站日常运营，做成标准、可复制的流程。")


def slide_pain(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_title_bar(slide, "您是否遇到过这些困扰？", "")
    pains = [
        ("上架慢", "每条链接都要传图、填属性、对价格，占用大量人力"),
        ("跟踪难", "新发产品多了，很难持续盯住每一条的表现"),
        ("数据散", "参谋、关键词等数据分散在多个表格，整理费时"),
        ("决策慢", "不清楚哪些该加推、哪些该优化，容易凭感觉操作"),
    ]
    y = 1.75
    for title, desc in pains:
        card = slide.shapes.add_shape(1, Inches(0.65), Inches(y), Inches(8.7), Inches(1.15))
        card.fill.solid()
        card.fill.fore_color.rgb = LIGHT_BG
        card.line.fill.background()
        tf = card.text_frame
        tf.margin_left = Inches(0.2)
        p0 = tf.paragraphs[0]
        p0.text = title
        p0.font.size = Pt(18)
        p0.font.bold = True
        p0.font.color.rgb = ACCENT
        p1 = tf.add_paragraph()
        p1.text = desc
        p1.font.size = Pt(14)
        p1.font.color.rgb = DARK
        p1.space_before = Pt(4)
        y += 1.28


def slide_solution(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_title_bar(slide, "系统如何帮助您", "四大能力，覆盖上新到决策")
    add_feature_cards(
        slide,
        [
            ("智能发品", "批量/定时上架，自动完成图片、属性、价格与发布"),
            ("素材管理", "图片规范整理，支持 AI 辅助生图与店铺素材采集"),
            ("数据获取", "自动采集参谋、关键词、店铺与流量等运营数据"),
            ("分析看板", "综合分析、异动监测、P4P 与推荐关注一目了然"),
        ],
    )


def slide_flow(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_title_bar(slide, "您的使用流程", "简单四步，无需复杂技术背景")
    steps = [
        ("① 准备资料", "按模板整理产品 Excel 与图片素材"),
        ("② 配置规则", "在界面中设置路径、属性与价格规则（可协助初始化）"),
        ("③ 一键执行", "启动发品或数据任务，进度实时可见"),
        ("④ 查看结果", "在分析页与控制台查看该优化、该跟进的商品"),
    ]
    x = 0.55
    for title, desc in steps:
        card = slide.shapes.add_shape(1, Inches(x), Inches(1.85), Inches(2.1), Inches(4.5))
        card.fill.solid()
        card.fill.fore_color.rgb = LIGHT_BG
        card.line.color.rgb = PRIMARY
        tf = card.text_frame
        tf.margin_top = Inches(0.3)
        p0 = tf.paragraphs[0]
        p0.text = title
        p0.font.size = Pt(17)
        p0.font.bold = True
        p0.font.color.rgb = PRIMARY
        p0.alignment = PP_ALIGN.CENTER
        p1 = tf.add_paragraph()
        p1.text = desc
        p1.font.size = Pt(12)
        p1.font.color.rgb = DARK
        p1.alignment = PP_ALIGN.CENTER
        p1.space_before = Pt(14)
        x += 2.35
    add_quote_bar(slide, "支持本机安装使用，数据保存在您的电脑，使用更安心。", top=6.55)


def slide_publish(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_title_bar(slide, "智能发品", "从表格到上线，自动完成")
    add_bullets(
        slide,
        [
            "自动完成：登录、读取产品表、属性填写、图片上传、价格设置、提交发布",
            "支持批量发布与定时发布，灵活匹配您的上新节奏",
            "发布成功后，新产品自动纳入「新发链接」跟踪列表",
            "可配合「产品优化」「视频绑定」等后续运营动作",
            "任务过程透明，随时查看执行进度与结果",
        ],
        top=1.7,
        font_size=17,
    )


def slide_newlink(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_title_bar(slide, "新发产品，自动跟踪", "不用再手工记录产品编号")
    add_bullets(
        slide,
        [
            "每一条成功发布的产品，系统会自动记录产品 ID 与发布日期",
            "结合店铺数据分析，按周曝光等维度筛选表现突出的新品",
            "在控制台「推荐关注」区域集中展示，方便优先跟进",
            "帮助您回答：哪些新品值得加橱窗、加推广、重点优化",
        ],
        top=1.85,
        font_size=18,
    )
    add_quote_bar(slide, "上新之后，也能看得见、跟得上。")


def slide_data(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_title_bar(slide, "运营数据，自动汇总", "减少每周手工下载与拼表")
    add_bullets(
        slide,
        [
            "可自动获取多类运营数据，例如：",
            "  · 产品参谋数据（曝光、点击、询盘等）",
            "  · 关键词与行业词数据",
            "  · 店铺运营与流量渠道数据",
            "一键运行综合分析，生成统计结果与监控报表",
            "用数据支撑每周运营会议，而不是凭印象开会",
        ],
        top=1.65,
        font_size=17,
    )


def slide_insight(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_title_bar(slide, "分析看板，辅助决策", "告诉您「接下来该做什么」")
    items = [
        ("综合分析", "多维度指标统计，了解全店产品表现"),
        ("产品诊断", "健康度评估，识别需要优化的链接"),
        ("P4P 分析", "区分有效果的推广与低效投放"),
        ("异动监测", "快速发现曝光、点击明显变化的产品"),
        ("推荐关注", "综合表现优秀 + 潜力新品，优先处理"),
    ]
    y = 1.75
    for title, desc in items:
        row = slide.shapes.add_shape(1, Inches(0.65), Inches(y), Inches(8.7), Inches(0.95))
        row.fill.solid()
        row.fill.fore_color.rgb = LIGHT_BG
        row.line.fill.background()
        tf = row.text_frame
        tf.margin_left = Inches(0.18)
        p0 = tf.paragraphs[0]
        p0.text = f"{title}  —  {desc}"
        p0.font.size = Pt(15)
        p0.font.color.rgb = DARK
        y += 1.05


def slide_dashboard(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_title_bar(slide, "运营控制台", "关键信息，一屏掌握")
    add_bullets(
        slide,
        [
            "店铺运营状态总览，快速进入诊断与优化",
            "P4P 投放效果摘要、产品排名与异动提醒",
            "推荐关注清单：该加推、该优化、该跟进的新品",
            "支持将发品、优化、数据下载、分析串联为完整工作流",
        ],
        top=1.85,
        font_size=18,
    )
    # 截图占位
    ph = slide.shapes.add_shape(1, Inches(5.8), Inches(2.0), Inches(3.5), Inches(3.8))
    ph.fill.solid()
    ph.fill.fore_color.rgb = RGBColor(0xEA, 0xEF, 0xF4)
    ph.line.color.rgb = SOFT
    ph.text_frame.text = "系统界面\n示意图"
    ph.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    ph.text_frame.paragraphs[0].font.size = Pt(14)
    ph.text_frame.paragraphs[0].font.color.rgb = SOFT


def slide_image(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_title_bar(slide, "图片与素材支持", "让上新前的准备更高效")
    add_bullets(
        slide,
        [
            "图片规范化：统一分组与命名，发品更顺畅",
            "AI 生图：根据配置批量生成主图、场景图（按需使用）",
            "店铺图片采集：整理已有素材，方便复用",
        ],
        top=2.2,
        font_size=19,
    )


def slide_suitable(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_title_bar(slide, "适合怎样的店铺？", "")
    add_feature_cards(
        slide,
        [
            ("SKU 较多", "需要稳定、高频上新的铺货型店铺"),
            ("重视数据", "希望用数据指导推广与优化决策"),
            ("团队精干", "希望减少重复岗位、提升人效"),
            ("多店运营", "需要标准化流程、可复制配置"),
        ],
        top=1.85,
    )


def slide_security(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_title_bar(slide, "数据与使用安全", "我们理解您对账号与数据的关注")
    add_bullets(
        slide,
        [
            "系统在您本机运行，运营数据主要保存在本地",
            "配置与任务过程可追溯，便于团队内部协作与交接",
            "支持会员账号管理，按授权使用功能",
            "请在平台规则范围内合规使用自动化功能",
            "我们提供安装指导与使用培训，帮助您快速上手",
        ],
        top=1.75,
        font_size=17,
    )


def slide_service(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_title_bar(slide, "我们提供的服务", "让您用得起来、用得好")
    add_bullets(
        slide,
        [
            "系统安装与基础配置协助",
            "首次发品陪跑，确保流程跑通",
            "使用培训（界面操作、日常运营节奏建议）",
            "问题响应与版本更新支持（以合同约定为准）",
        ],
        top=2.0,
        font_size=19,
    )


def slide_start(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_title_bar(slide, "如何开始？", "")
    steps = [
        ("了解演示", "预约在线或现场演示，直观感受系统界面与流程"),
        ("试用体验", "在试用期内体验核心发品与分析功能"),
        ("正式使用", "确定方案后完成安装配置，进入日常运营"),
    ]
    y = 2.0
    for i, (title, desc) in enumerate(steps, 1):
        circle = slide.shapes.add_shape(1, Inches(0.75), Inches(y), Inches(0.55), Inches(0.55))
        circle.fill.solid()
        circle.fill.fore_color.rgb = ACCENT
        circle.line.fill.background()
        circle.text_frame.text = str(i)
        circle.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        circle.text_frame.paragraphs[0].font.size = Pt(18)
        circle.text_frame.paragraphs[0].font.bold = True
        circle.text_frame.paragraphs[0].font.color.rgb = WHITE

        tb = slide.shapes.add_textbox(Inches(1.5), Inches(y - 0.05), Inches(7.8), Inches(1.0))
        tf = tb.text_frame
        tf.text = f"{title}\n{desc}"
        tf.paragraphs[0].font.size = Pt(18)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = PRIMARY
        tf.paragraphs[1].font.size = Pt(14)
        tf.paragraphs[1].font.color.rgb = DARK
        y += 1.45


def slide_contact(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, PRIMARY)
    t1 = slide.shapes.add_textbox(Inches(0.8), Inches(2.2), Inches(8.4), Inches(0.9))
    t1.text_frame.text = "欢迎进一步了解"
    t1.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    t1.text_frame.paragraphs[0].font.size = Pt(36)
    t1.text_frame.paragraphs[0].font.bold = True
    t1.text_frame.paragraphs[0].font.color.rgb = WHITE

    info = slide.shapes.add_textbox(Inches(1.2), Inches(3.3), Inches(7.6), Inches(2.5))
    tf = info.text_frame
    tf.text = (
        "预约产品演示\n"
        "申请试用体验\n\n"
        "服务热线：________________\n"
        "微信 / 邮箱：________________\n"
        "官方网站：________________"
    )
    for i, p in enumerate(tf.paragraphs):
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(20 if i < 2 else 16)
        p.font.color.rgb = RGBColor(0xD8, 0xE8, 0xF8)
        if i < 2:
            p.font.bold = True

    t3 = slide.shapes.add_textbox(Inches(0.8), Inches(6.0), Inches(8.4), Inches(0.5))
    t3.text_frame.text = "感谢您的关注"
    t3.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    t3.text_frame.paragraphs[0].font.size = Pt(14)
    t3.text_frame.paragraphs[0].font.color.rgb = RGBColor(0x98, 0xB8, 0xD0)


def slide_end(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, PRIMARY)
    t1 = slide.shapes.add_textbox(Inches(1), Inches(2.9), Inches(8), Inches(1.2))
    tf = t1.text_frame
    tf.text = "专注国际站运营效率\n助力您的生意持续增长"
    for p in tf.paragraphs:
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(30)
        p.font.bold = True
        p.font.color.rgb = WHITE


def build():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    slide_cover(prs)
    slide_about(prs)
    slide_pain(prs)
    slide_solution(prs)
    slide_flow(prs)
    slide_publish(prs)
    slide_newlink(prs)
    slide_data(prs)
    slide_insight(prs)
    slide_dashboard(prs)
    slide_image(prs)
    slide_suitable(prs)
    slide_security(prs)
    slide_service(prs)
    slide_start(prs)
    slide_contact(prs)
    slide_end(prs)

    OUTPUT_CN.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUTPUT_CN))
    prs.save(str(OUTPUT_EN))
    print(f"已生成: {OUTPUT_CN}")
    print(f"已生成: {OUTPUT_EN}")
    print(f"共 {len(prs.slides)} 页")


if __name__ == "__main__":
    build()
