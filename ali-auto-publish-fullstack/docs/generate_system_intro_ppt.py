# -*- coding: utf-8 -*-
"""生成「系统介绍」PPT。运行: py -3.11 docs/generate_system_intro_ppt.py"""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parent.parent
OUTPUT = ROOT / "docs" / "阿里巴巴国际站-自动发品与运营管理系统-介绍.pptx"

# 品牌色
PRIMARY = RGBColor(0x0F, 0x4C, 0x81)  # 深蓝
ACCENT = RGBColor(0xFF, 0x6A, 0x00)  # 阿里橙点缀
LIGHT_BG = RGBColor(0xF0, 0xF4, 0xF8)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK = RGBColor(0x1A, 0x1A, 0x2E)
MUTED = RGBColor(0x5A, 0x6A, 0x7A)


def set_slide_bg(slide, rgb: RGBColor):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = rgb


def add_title_bar(slide, title: str, subtitle: str = ""):
    """顶部标题条"""
    bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(1.05))
    bar.fill.solid()
    bar.fill.fore_color.rgb = PRIMARY
    bar.line.fill.background()
    tf = bar.text_frame
    tf.text = title
    p = tf.paragraphs[0]
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = WHITE
    if subtitle:
        box = slide.shapes.add_textbox(Inches(0.6), Inches(1.15), Inches(8.8), Inches(0.45))
        tfb = box.text_frame
        tfb.text = subtitle
        tfb.paragraphs[0].font.size = Pt(14)
        tfb.paragraphs[0].font.color.rgb = MUTED


def add_bullets(slide, items: list[str], left=0.7, top=1.75, width=8.6, height=5.0, font_size=18):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    for i, text in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = text
        p.level = 0
        p.font.size = Pt(font_size)
        p.font.color.rgb = DARK
        p.space_after = Pt(10)
        if text.startswith("  "):
            p.level = 1
            p.font.size = Pt(font_size - 2)


def add_two_column(slide, left_title, left_items, right_title, right_items):
    add_bullets(slide, [left_title, ""] + left_items, left=0.6, top=1.7, width=4.2, font_size=16)
    add_bullets(slide, [right_title, ""] + right_items, left=5.2, top=1.7, width=4.2, font_size=16)


def slide_cover(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, PRIMARY)
    t1 = slide.shapes.add_textbox(Inches(0.8), Inches(2.2), Inches(8.4), Inches(1.2))
    tf = t1.text_frame
    tf.text = "阿里巴巴国际站"
    tf.paragraphs[0].font.size = Pt(22)
    tf.paragraphs[0].font.color.rgb = RGBColor(0xB8, 0xD4, 0xE8)

    t2 = slide.shapes.add_textbox(Inches(0.8), Inches(2.9), Inches(8.4), Inches(1.4))
    tf2 = t2.text_frame
    tf2.text = "自动发品与运营管理系统"
    tf2.paragraphs[0].font.size = Pt(40)
    tf2.paragraphs[0].font.bold = True
    tf2.paragraphs[0].font.color.rgb = WHITE

    t3 = slide.shapes.add_textbox(Inches(0.8), Inches(4.5), Inches(8.4), Inches(0.8))
    tf3 = t3.text_frame
    tf3.text = "发品自动化 · 数据可视化 · 运营有章法"
    tf3.paragraphs[0].font.size = Pt(20)
    tf3.paragraphs[0].font.color.rgb = RGBColor(0xE8, 0xF0, 0xF8)

    line = slide.shapes.add_shape(1, Inches(0.8), Inches(5.4), Inches(2.2), Inches(0.06))
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT
    line.line.fill.background()

    t4 = slide.shapes.add_textbox(Inches(0.8), Inches(5.6), Inches(8), Inches(0.5))
    tf4 = t4.text_frame
    tf4.text = "全栈版 · FastAPI + React · 本机/桌面部署"
    tf4.paragraphs[0].font.size = Pt(14)
    tf4.paragraphs[0].font.color.rgb = RGBColor(0xA0, 0xB8, 0xCC)


def slide_agenda(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_title_bar(slide, "目录", "Agenda")
    add_bullets(
        slide,
        [
            "01  行业痛点与产品定位",
            "02  系统架构与技术栈",
            "03  四大业务模块总览",
            "04  自动发品：七步流程",
            "05  图片管理与 AI 生图",
            "06  数据自动下载",
            "07  数据分析与决策看板",
            "08  控制台与运营闭环",
            "09  核心优势与适用对象",
        ],
        top=1.85,
        font_size=20,
    )


def slide_pain(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_title_bar(slide, "行业痛点", "国际站卖家日常运营中的典型挑战")
    add_two_column(
        slide,
        "常见痛点",
        [
            "• 后台逐条发品，属性/价格/图片重复劳动",
            "• Excel 与平台来回切换，易漏填、易错价",
            "• 新发产品 ID 手工记录，跟踪困难",
            "• 参谋/关键词数据每周手工导出",
            "• 靠经验判断优化对象，缺乏统一指标",
            "• 多个脚本分散，难维护、难协作",
        ],
        "我们的定位",
        [
            "• 专为 Alibaba.com 定制的运营工具",
            "• 发品 + 优化 + 绑视频 一体化",
            "• 数据下载 → 分析 → 看板 闭环",
            "• 可视化配置，规则可复用",
            "• 实时任务日志，过程可追溯",
            "• 支持桌面安装包，团队易分发",
        ],
    )


def slide_arch(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_title_bar(slide, "系统架构", "前后端分离 · 模块化服务层 · Monorepo")
    add_bullets(
        slide,
        [
            "【前端】React 19 + Vite + TailwindCSS",
            "  → 可视化配置、任务控制、数据分析看板、WebSocket 实时日志",
            "",
            "【后端】FastAPI + Uvicorn + Pydantic",
            "  → API 路由层 + Service 业务层（发品自动化 / 图片 / 下载 / 分析）",
            "",
            "【自动化引擎】Playwright / Selenium + Chrome/Edge",
            "  → 登录、翻页下载、填表发品、图片上传",
            "",
            "【数据存储】本地 Excel + 可配置路径（数据不出本机，便于私有化部署）",
            "",
            "【交付形态】开发模式（前后端分离启动）/ 桌面安装包（PyInstaller + Electron）",
        ],
        top=1.65,
        font_size=17,
    )


def slide_modules(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_title_bar(slide, "四大业务模块", "覆盖「上架 — 素材 — 数据 — 决策」全链路")
    cards = [
        ("产品上传", "自动发品 · 优化产品 · 绑视频 · 配置管理"),
        ("图片管理", "图片规范化 · AI 生图 · 店铺图片采集"),
        ("数据下载", "参谋 · 运营 · 关键词 · 行业词 · 店铺 · 流量渠道"),
        ("数据分析", "综合分析 · 周诊断 · P4P · 新发监控 · 单品/流量分析"),
    ]
    positions = [(0.5, 1.8), (5.1, 1.8), (0.5, 4.0), (5.1, 4.0)]
    for (title, desc), (x, y) in zip(cards, positions):
        shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(4.4), Inches(1.85))
        shape.fill.solid()
        shape.fill.fore_color.rgb = LIGHT_BG
        shape.line.color.rgb = PRIMARY
        tf = shape.text_frame
        tf.margin_left = Inches(0.15)
        tf.margin_top = Inches(0.12)
        p0 = tf.paragraphs[0]
        p0.text = title
        p0.font.size = Pt(22)
        p0.font.bold = True
        p0.font.color.rgb = PRIMARY
        p1 = tf.add_paragraph()
        p1.text = desc
        p1.font.size = Pt(13)
        p1.font.color.rgb = DARK
        p1.space_before = Pt(8)


def slide_publish(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_title_bar(slide, "自动发品 · 七步流程", "表格准备好，其余交给自动化")
    steps = [
        "1. 登录验证 — Cookie 或浏览器登录国际站",
        "2. 读取数据源 — 从 Excel 读取产品与配置",
        "3. 属性融合 — 合并属性、自定义映射规则",
        "4. 图片上传 — 首图、详情图、SKU 图",
        "5. 价格设置 — 阶梯价格、FOB 价格",
        "6. 填写属性 — 自动填写规格与类目属性",
        "7. 发布产品 — 提交并校验成功",
    ]
    add_bullets(slide, steps, top=1.7, font_size=17)
    foot = slide.shapes.add_textbox(Inches(0.7), Inches(5.85), Inches(8.6), Inches(0.55))
    ftf = foot.text_frame
    ftf.text = (
        "扩展能力：批量/定时/每日定时发品 · 橱窗/P4P 开关 · 发品成功自动写入「新发链接监控」"
    )
    ftf.paragraphs[0].font.size = Pt(13)
    ftf.paragraphs[0].font.color.rgb = ACCENT
    ftf.paragraphs[0].font.bold = True


def slide_image(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_title_bar(slide, "图片管理", "素材效率与规范一并提升")
    add_bullets(
        slide,
        [
            "图片规范化",
            "  • 按分组管理素材，场景识别与命名规范",
            "  • 批量标准化处理，减少发品前手工整理",
            "",
            "AI 生图",
            "  • 可配置提示词、输入输出目录",
            "  • 积分预估与任务启停，批量生成主图/场景图",
            "",
            "店铺图片采集",
            "  • 从店铺侧采集图片资源，便于复用与二次编辑",
        ],
        top=1.65,
        font_size=17,
    )


def slide_download(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_title_bar(slide, "数据自动下载", "浏览器自动化翻页采集，支持任务监控")
    add_bullets(
        slide,
        [
            "产品参谋数据 — 曝光、点击、询盘等核心指标",
            "产品运营数据 — 运营维度明细表",
            "关键词数据 — 词包表现与异动监测",
            "行业关键词 — 行业词库维护，支持标题生成任务",
            "店铺运营数据 — 店铺级运营概览",
            "流量渠道数据 — 自然/场景/推广等渠道拆分",
            "",
            "任务能力：启动 / 停止 / 状态查询 · 增量更新 · 与综合分析任务联动",
        ],
        top=1.65,
        font_size=17,
    )


def slide_analysis(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_title_bar(slide, "数据分析", "从原始 Excel 到可执行决策")
    add_bullets(
        slide,
        [
            "综合分析 — 多指标统计（曝光/点击/访问/收藏/询盘/TM 等）",
            "  权重评分、推荐关注、新发链接监控表自动生成",
            "",
            "周数据分析 / 产品诊断 — 健康度分层、连续评分、阈值可配置",
            "P4P 分析 — 有询盘产品 vs 低点击无询盘产品",
            "新发链接监控 — 按周曝光筛选高潜力新品",
            "产品优化建议 — 标题等优化方向分析",
            "单品分析 · 单品渠道 · 流量分析 — SKU 与渠道下沉",
        ],
        top=1.65,
        font_size=16,
    )


def slide_dashboard(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_title_bar(slide, "控制台 · 运营驾驶舱", "一屏掌握店铺关键信号")
    add_two_column(
        slide,
        "看板模块",
        [
            "• 店铺诊断入口",
            "• P4P 分析摘要",
            "• 产品综合排名",
            "• 产品异动（涨跌最明显）",
            "• 推荐关注 + 新发监控",
        ],
        "运营闭环（可一键串联）",
        [
            "① 自动发品",
            "② 优化产品",
            "③ 新品绑定视频",
            "④ 数据下载（多类型）",
            "⑤ 综合分析",
            "→ 形成「上架—跟踪—优化」闭环",
        ],
    )


def slide_flow_newlink(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_title_bar(slide, "新发链接自动追踪", "发品成功即入库，分析自动过滤")
    add_bullets(
        slide,
        [
            "发品成功后，系统从成功 URL 解析 primaryId（产品 ID）",
            "自动写入「新发链接监控」源表（发品日期 + 新发链接，去重、置顶）",
            "综合分析任务读取新发名单 + 统计指标，生成「新发链接数据监控.xlsx」",
            "前端按周曝光规则筛选（如近 4 周中至少 3 周曝光 > 30）",
            "控制台「推荐关注产品」区域同步展示高潜力新品",
            "",
            "价值：无需手工记 ID，新品表现可量化、可对比、可优先跟进",
        ],
        top=1.65,
        font_size=17,
    )


def slide_config(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_title_bar(slide, "可视化配置与实时监控", "降低使用门槛，提升团队协作效率")
    add_bullets(
        slide,
        [
            "配置管理（前端界面）",
            "  • 文件路径、属性映射、阶梯价格公式、分组 URL",
            "  • 从平台拉取属性/规格，Cookie 与浏览器登录管理",
            "",
            "实时任务监控",
            "  • WebSocket 推送任务日志，发品/下载/分析进度实时可见",
            "  • 支持暂停、停止、定时任务",
            "",
            "会员中心",
            "  • 注册登录、试用与会员策略，适合授权/分发商业模式",
        ],
        top=1.65,
        font_size=17,
    )


def slide_tech(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_title_bar(slide, "技术栈与扩展性", "架构为二次开发与长期维护而设计")
    add_two_column(
        slide,
        "技术栈",
        [
            "后端：FastAPI, Uvicorn, Pydantic",
            "前端：React 19, Vite, TailwindCSS 4",
            "UI：Radix UI, Lucide 图标",
            "自动化：Playwright / Selenium",
            "实时：WebSocket 日志与任务",
        ],
        "扩展指南（示例）",
        [
            "改属性映射 → attribute_filler.py",
            "改阶梯价格 → price_setter.py",
            "改页面元素 → page_helpers.py",
            "新增配置项 → settings.py + ProductConfig 页",
            "新增 API → api 路由 + 前端 api.ts",
        ],
    )


def slide_deploy(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_title_bar(slide, "部署与交付", "开发环境 + 客户桌面包双模式")
    add_bullets(
        slide,
        [
            "开发/内网部署",
            "  • 后端：Python 3.9+，pip install，python run.py（默认 :8000）",
            "  • 前端：Node 18+，pnpm install，pnpm run dev（默认 :3000）",
            "",
            "客户分发（桌面安装包）",
            "  • PyInstaller 打包后端 exe",
            "  • Electron 桌面壳 + 前端静态资源",
            "  • 本机运行，数据本地存储，适合私有化交付",
        ],
        top=1.65,
        font_size=17,
    )


def slide_value(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_title_bar(slide, "核心优势总结", "为什么选择本系统")
    items = [
        ("效率", "7 步自动发品，批量+定时，显著减少重复操作"),
        ("闭环", "发品 → 监控 → 下载 → 分析 → 看板，数据驱动运营"),
        ("可控", "可视化配置 + 实时日志，规则透明、问题可追溯"),
        ("专业", "面向国际站场景，覆盖 P4P、新发、异动、排名等维度"),
        ("可交付", "模块化代码 + 桌面安装包，便于定制与商业分发"),
    ]
    y = 1.75
    for title, desc in items:
        dot = slide.shapes.add_shape(1, Inches(0.65), Inches(y + 0.08), Inches(0.18), Inches(0.18))
        dot.fill.solid()
        dot.fill.fore_color.rgb = ACCENT
        dot.line.fill.background()
        tb = slide.shapes.add_textbox(Inches(1.0), Inches(y), Inches(8.3), Inches(0.75))
        tf = tb.text_frame
        tf.text = f"{title}：{desc}"
        tf.paragraphs[0].font.size = Pt(17)
        tf.paragraphs[0].font.color.rgb = DARK
        y += 0.82


def slide_audience(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_title_bar(slide, "适用对象", "Who is it for?")
    add_bullets(
        slide,
        [
            "阿里巴巴国际站 B2B 卖家 — 多 SKU、铺货型、需稳定上新",
            "运营团队 — 需要标准化发品流程与统一数据口径",
            "代运营公司 — 多店铺可复制配置，桌面包便于客户交付",
            "技术型卖家 — 希望在本机私有化部署、可二次开发扩展",
        ],
        top=2.0,
        font_size=20,
    )


def slide_end(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, PRIMARY)
    t1 = slide.shapes.add_textbox(Inches(1), Inches(2.6), Inches(8), Inches(1))
    tf = t1.text_frame
    tf.text = "感谢聆听"
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].font.size = Pt(44)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = WHITE

    t2 = slide.shapes.add_textbox(Inches(1), Inches(3.8), Inches(8), Inches(1.2))
    tf2 = t2.text_frame
    tf2.text = "阿里巴巴国际站 · 自动发品与运营管理系统\n\n欢迎演示交流 · Q&A"
    for p in tf2.paragraphs:
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(0xC8, 0xDC, 0xEC)


def build():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    slide_cover(prs)
    slide_agenda(prs)
    slide_pain(prs)
    slide_arch(prs)
    slide_modules(prs)
    slide_publish(prs)
    slide_image(prs)
    slide_download(prs)
    slide_analysis(prs)
    slide_flow_newlink(prs)
    slide_dashboard(prs)
    slide_config(prs)
    slide_tech(prs)
    slide_deploy(prs)
    slide_value(prs)
    slide_audience(prs)
    slide_end(prs)

    OUTPUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUTPUT))
    print(f"已生成: {OUTPUT}")
    print(f"共 {len(prs.slides)} 页")


if __name__ == "__main__":
    build()
