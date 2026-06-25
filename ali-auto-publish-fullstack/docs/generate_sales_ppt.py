# -*- coding: utf-8 -*-
"""生成「销售端」PPT。运行: py -3.11 docs/generate_sales_ppt.py"""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parent.parent
OUTPUT_CN = ROOT / "docs" / "销售端-国际站运营自动化方案.pptx"
OUTPUT_EN = ROOT / "docs" / "sales-deck.pptx"

PRIMARY = RGBColor(0x0F, 0x4C, 0x81)
ACCENT = RGBColor(0xFF, 0x6A, 0x00)
GOLD = RGBColor(0xC9, 0x8A, 0x1A)
LIGHT_BG = RGBColor(0xF8, 0xFA, 0xFC)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK = RGBColor(0x1A, 0x1A, 0x2E)
MUTED = RGBColor(0x5A, 0x6A, 0x7A)
GREEN = RGBColor(0x16, 0x8A, 0x4A)
RED = RGBColor(0xC0, 0x39, 0x2B)


def set_slide_bg(slide, rgb: RGBColor):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = rgb


def add_title_bar(slide, title: str, subtitle: str = ""):
    bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(1.05))
    bar.fill.solid()
    bar.fill.fore_color.rgb = PRIMARY
    bar.line.fill.background()
    tf = bar.text_frame
    tf.text = title
    p = tf.paragraphs[0]
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = WHITE
    if subtitle:
        box = slide.shapes.add_textbox(Inches(0.6), Inches(1.12), Inches(8.8), Inches(0.45))
        tfb = box.text_frame
        tfb.text = subtitle
        tfb.paragraphs[0].font.size = Pt(14)
        tfb.paragraphs[0].font.color.rgb = MUTED


def add_bullets(slide, items, left=0.7, top=1.75, width=8.6, font_size=18, color=DARK):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(5.2))
    tf = box.text_frame
    tf.word_wrap = True
    for i, text in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = text
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.space_after = Pt(8)
        if text.startswith("  "):
            p.level = 1
            p.font.size = Pt(font_size - 2)


def add_highlight_box(slide, text: str, top=5.9, color=ACCENT):
    box = slide.shapes.add_textbox(Inches(0.65), Inches(top), Inches(8.7), Inches(0.55))
    tf = box.text_frame
    tf.text = text
    p = tf.paragraphs[0]
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = color


def add_two_column(slide, left_title, left_items, right_title, right_items, top=1.7):
    for title, items, left in [(left_title, left_items, 0.55), (right_title, right_items, 5.15)]:
        hdr = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(4.3), Inches(0.4))
        hdr.text_frame.text = title
        hdr.text_frame.paragraphs[0].font.size = Pt(18)
        hdr.text_frame.paragraphs[0].font.bold = True
        hdr.text_frame.paragraphs[0].font.color.rgb = PRIMARY
        add_bullets(slide, items, left=left, top=top + 0.45, width=4.25, font_size=15)


def add_stat_row(slide, stats: list[tuple[str, str, str]], top=2.0):
    """三列数据：数字 / 标签 / 说明"""
    w = 2.9
    for i, (num, label, note) in enumerate(stats):
        x = 0.55 + i * 3.15
        card = slide.shapes.add_shape(1, Inches(x), Inches(top), Inches(w), Inches(2.6))
        card.fill.solid()
        card.fill.fore_color.rgb = LIGHT_BG
        card.line.color.rgb = PRIMARY
        tf = card.text_frame
        tf.margin_top = Inches(0.2)
        p0 = tf.paragraphs[0]
        p0.text = num
        p0.font.size = Pt(32)
        p0.font.bold = True
        p0.font.color.rgb = ACCENT
        p0.alignment = PP_ALIGN.CENTER
        p1 = tf.add_paragraph()
        p1.text = label
        p1.font.size = Pt(16)
        p1.font.bold = True
        p1.font.color.rgb = DARK
        p1.alignment = PP_ALIGN.CENTER
        p1.space_before = Pt(6)
        p2 = tf.add_paragraph()
        p2.text = note
        p2.font.size = Pt(11)
        p2.font.color.rgb = MUTED
        p2.alignment = PP_ALIGN.CENTER
        p2.space_before = Pt(4)


def slide_cover(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, PRIMARY)
    t0 = slide.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(8.4), Inches(0.5))
    t0.text_frame.text = "客户方案 · 销售专用"
    t0.text_frame.paragraphs[0].font.size = Pt(14)
    t0.text_frame.paragraphs[0].font.color.rgb = RGBColor(0xA8, 0xC8, 0xE0)

    t1 = slide.shapes.add_textbox(Inches(0.8), Inches(2.1), Inches(8.4), Inches(1.6))
    tf = t1.text_frame
    tf.text = "国际站运营，别再靠人海战术"
    tf.paragraphs[0].font.size = Pt(38)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = WHITE

    t2 = slide.shapes.add_textbox(Inches(0.8), Inches(3.5), Inches(8.4), Inches(1.0))
    tf2 = t2.text_frame
    tf2.text = "自动发品与运营管理系统\n一条链路搞定：上架 · 跟踪 · 下载 · 分析 · 决策"
    for p in tf2.paragraphs:
        p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(0xD8, 0xE8, 0xF4)

    line = slide.shapes.add_shape(1, Inches(0.8), Inches(5.0), Inches(2.5), Inches(0.07))
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT
    line.line.fill.background()

    t3 = slide.shapes.add_textbox(Inches(0.8), Inches(5.25), Inches(8.2), Inches(0.9))
    tf3 = t3.text_frame
    tf3.text = "【可填写】公司名称  |  销售顾问  |  演示日期"
    tf3.paragraphs[0].font.size = Pt(13)
    tf3.paragraphs[0].font.color.rgb = RGBColor(0x90, 0xB0, 0xC8)


def slide_hook(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_title_bar(slide, "您是否也在经历这些？", "典型国际站卖家的日常")
    add_bullets(
        slide,
        [
            "运营每天泡在后台：传图、填属性、对价格，一条链接半小时起",
            "店铺 SKU 一多，新人培训成本高，老人一走流程就断档",
            "新发的产品记不全 ID，上了架也不知道后面曝光好不好",
            "参谋数据要人工导，周报靠 Excel 拼，决策总是慢半拍",
            "P4P 在烧钱，但说不清哪些品该加推、哪些该止损",
            "手里好几个脚本，出问题不知道卡在哪一步",
        ],
        top=1.7,
        font_size=17,
    )
    add_highlight_box(slide, "销售话术：不是您不努力，是重复劳动把团队精力耗光了。")


def slide_solution(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_title_bar(slide, "我们提供的答案", "一套系统，把国际站运营做成「标准流水线」")
    add_stat_row(
        slide,
        [
            ("省时", "自动发品", "7步流程批量/定时上架"),
            ("省心", "自动跟品", "新发ID自动入库监控"),
            ("省力", "数据决策", "下载+分析+看板一体"),
        ],
        top=2.1,
    )
    add_bullets(
        slide,
        [
            "您只需：按模板准备好 Excel 与图片，配置好规则",
            "系统完成：登录、融属性、传图、设价、发布、记新发、拉数据、出分析",
            "您得到：该上新的上新，该优化的优化，该加推的加推——有名单、有依据",
        ],
        top=4.95,
        font_size=16,
    )


def slide_before_after(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_title_bar(slide, "使用前 vs 使用后", "同样团队规模，做更多事、做错更少")
    add_two_column(
        slide,
        "使用前（人工）",
        [
            "• 1 人每天发 5～10 条已属高效",
            "• 属性/价格靠记忆，易出错被平台打回",
            "• 新发品跟踪靠表格手工维护",
            "• 周报 = 下载 + 复制 + 透视",
            "• 优化靠经验，难以复盘",
        ],
        "使用后（系统）",
        [
            "• 批量+定时，夜间跑量白天盯询盘",
            "• 规则一次配置，全店统一执行",
            "• 发品成功自动写入监控池",
            "• 一键下载 + 综合分析出结果",
            "• 控制台看异动、P4P、推荐名单",
        ],
    )
    add_highlight_box(
        slide,
        "【销售填写】客户实测：发品效率提升 ____ 倍  |  运营节省 ____ 人天/周",
        color=GOLD,
    )


def slide_value_pillars(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_title_bar(slide, "三大客户价值", "对客户只说结果，不说技术名词")
    cards = [
        ("多卖货", "稳定上新 + 新发监控", "持续铺货不中断，高曝光新品及时跟进转化"),
        ("少烧钱", "P4P/异动/推荐看板", "把钱花在有点击、有询盘的品上，低效品早识别"),
        ("少招人", "流程标准化", "1 套系统顶掉大量重复岗位，团队聚焦选品与谈单"),
    ]
    y = 1.85
    for title, sub, desc in cards:
        shape = slide.shapes.add_shape(1, Inches(0.55), Inches(y), Inches(8.9), Inches(1.55))
        shape.fill.solid()
        shape.fill.fore_color.rgb = LIGHT_BG
        shape.line.color.rgb = PRIMARY
        tf = shape.text_frame
        tf.margin_left = Inches(0.2)
        p0 = tf.paragraphs[0]
        p0.text = f"{title}  |  {sub}"
        p0.font.size = Pt(20)
        p0.font.bold = True
        p0.font.color.rgb = PRIMARY
        p1 = tf.add_paragraph()
        p1.text = desc
        p1.font.size = Pt(14)
        p1.font.color.rgb = DARK
        p1.space_before = Pt(6)
        y += 1.75


def slide_product_publish(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_title_bar(slide, "能力一：自动发品中心", "让客户「听得懂」的功能表述")
    add_bullets(
        slide,
        [
            "批量上架 — 表格导入，自动完成传图、属性、价格、发布",
            "定时上架 — 支持定时/每日定时，错峰发布更自然",
            "发完即建档 — 成功链接自动进入「新发监控」，不用手工记 ID",
            "后续自动化 — 可串联「优化产品」「新品绑视频」",
            "全程可查 — 任务日志实时可见，出问题立刻知道卡在哪一步",
        ],
        top=1.7,
        font_size=18,
    )
    add_highlight_box(slide, "成交点：把「发品专员」从重复劳动里解放出来。")


def slide_product_data(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_title_bar(slide, "能力二：数据自动采集 + 智能分析", "参谋数据不再靠人工周更")
    add_bullets(
        slide,
        [
            "自动下载 — 产品参谋、关键词、店铺运营、流量渠道等",
            "综合分析 — 曝光/点击/询盘/TM 等多维指标一键统计",
            "决策名单 — P4P 分析、产品异动、综合排名、推荐关注",
            "新发监控 — 按周曝光筛选潜力新品，优先跟进",
            "优化建议 — 标题等方向分析，减少「凭感觉改」",
        ],
        top=1.7,
        font_size=17,
    )
    add_highlight_box(slide, "成交点：老板每周看一张「该干什么」清单，而不是一堆原始表。")


def slide_product_assets(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_title_bar(slide, "能力三：图片与素材效率", "上新前的素材准备也能提速")
    add_bullets(
        slide,
        [
            "图片规范化 — 分组管理、命名统一，发品前少折腾",
            "AI 生图 — 批量生成主图/场景图，缩短拍摄与修图周期",
            "店铺图片采集 — 已有素材一键整理复用",
        ],
        top=2.0,
        font_size=19,
    )
    add_highlight_box(slide, "成交点：缩短「有想法」到「能上架」的时间。")


def slide_dashboard(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_title_bar(slide, "能力四：运营驾驶舱", "一屏回答：今天该盯哪些品？")
    add_bullets(
        slide,
        [
            "店铺诊断 — 快速进入问题与方案",
            "P4P 摘要 — 有询盘 vs 低点击无询盘，投放调整有依据",
            "异动榜单 — 涨跌最明显的 SKU 一眼可见",
            "推荐关注 — 综合评分优秀 + 新发高曝光潜力品",
            "一键串联 — 发品 → 优化 → 绑视频 → 下载 → 分析，全流程可自动化",
        ],
        top=1.7,
        font_size=17,
    )


def slide_scenarios(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_title_bar(slide, "典型成交场景", "对不同客户，强调不同痛点")
    scenarios = [
        ("铺货型卖家", "SKU 多、上新频", "强调批量发品 + 定时 + 新发监控"),
        ("成长型店铺", "想放大询盘", "强调数据分析 + P4P + 推荐关注"),
        ("代运营公司", "多店铺可复制", "强调标准流程 + 桌面包交付 + 会员授权"),
        ("老板亲自盯盘", "人少事多", "强调控制台一屏决策 + 省人力"),
    ]
    y = 1.75
    for who, pain, pitch in scenarios:
        row = slide.shapes.add_shape(1, Inches(0.55), Inches(y), Inches(8.9), Inches(1.2))
        row.fill.solid()
        row.fill.fore_color.rgb = LIGHT_BG
        row.line.fill.background()
        tf = row.text_frame
        tf.margin_left = Inches(0.15)
        p0 = tf.paragraphs[0]
        p0.text = who
        p0.font.size = Pt(17)
        p0.font.bold = True
        p0.font.color.rgb = ACCENT
        p1 = tf.add_paragraph()
        p1.text = f"痛点：{pain}    →    话术：{pitch}"
        p1.font.size = Pt(13)
        p1.font.color.rgb = DARK
        y += 1.35


def slide_compare(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_title_bar(slide, "为什么选择我们？", "对比表可直接用于谈判")
    rows = [
        ("对比项", "纯人工/零散脚本", "本系统"),
        ("发品效率", "低、依赖个人", "批量+定时+规则复用"),
        ("出错率", "高，难复盘", "流程固定，日志可追溯"),
        ("新发跟踪", "手工记 ID", "发品成功自动入库"),
        ("数据更新", "每周手工导", "自动下载+分析"),
        ("团队交接", "靠人带人", "配置即文档"),
        ("交付方式", "难复制", "支持桌面安装包+会员体系"),
        ("数据安全", "分散在多表", "本地存储，可私有化部署"),
    ]
    top = 1.65
    row_h = 0.52
    col_w = [1.5, 3.5, 3.5]
    col_x = [0.5, 2.1, 5.7]
    for r, row in enumerate(rows):
        y = top + r * row_h
        bg = PRIMARY if r == 0 else (LIGHT_BG if r % 2 == 0 else WHITE)
        for c, cell in enumerate(row):
            box = slide.shapes.add_shape(1, Inches(col_x[c]), Inches(y), Inches(col_w[c]), Inches(row_h))
            box.fill.solid()
            box.fill.fore_color.rgb = bg
            box.line.color.rgb = RGBColor(0xD0, 0xD8, 0xE0)
            tf = box.text_frame
            tf.margin_left = Inches(0.08)
            tf.vertical_anchor = 1
            p = tf.paragraphs[0]
            p.text = cell
            p.font.size = Pt(12 if r > 0 else 13)
            p.font.bold = r == 0 or c == 2
            p.font.color.rgb = WHITE if r == 0 else (PRIMARY if c == 2 and r > 0 else DARK)


def slide_trust(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_title_bar(slide, "客户最关心的四个问题", "售前异议处理")
    add_bullets(
        slide,
        [
            "Q：会不会很难学？",
            "  A：可视化配置，销售可演示「填表→一键发品」；核心操作 1 天内可上手。",
            "",
            "Q：阿里改版怎么办？",
            "  A：流程模块化维护，元素定位可更新；我们提供持续升级服务（按合同约定）。",
            "",
            "Q：数据会不会泄露？",
            "  A：默认本机运行、本地 Excel 存储，适合对数据敏感的客户。",
            "",
            "Q：适合几个人用？",
            "  A：小团队 1～2 人即可跑全店；代运营可多店铺复制配置。",
        ],
        top=1.6,
        font_size=15,
    )


def slide_delivery(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_title_bar(slide, "交付与服务", "让客户买得放心")
    add_bullets(
        slide,
        [
            "交付形态",
            "  • 桌面安装包 — 适合终端客户，双击即用",
            "  • 本机部署 — 适合对数据与网络有要求的团队",
            "",
            "会员体系（系统内置）",
            "  • 注册登录、试用期、VIP、充值与积分等能力",
            "  • 支持代理/分发商业模式（可按贵司政策包装）",
            "",
            "建议服务包（可写入合同）",
            "  • 远程安装与配置初始化",
            "  • 首次发品陪跑 + 周报模板对接",
            "  • 平台页面变更时的维护升级",
        ],
        top=1.65,
        font_size=16,
    )


def slide_pricing(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_title_bar(slide, "合作方案（请销售填写）", "报价页务必按公司最新政策更新")
    plans = [
        ("体验版", "试用 ____ 天", "核心发品 + 基础下载", "【价格】____ 元"),
        ("标准版", "单店年度", "发品+分析+看板全功能", "【价格】____ 元/年"),
        ("旗舰版", "单店+陪跑", "含部署、培训、季度优化", "【价格】____ 元/年"),
        ("代理/OEM", "多店铺分发", "桌面包+授权码+分润", "【政策】面议"),
    ]
    y = 1.75
    for name, term, scope, price in plans:
        card = slide.shapes.add_shape(1, Inches(0.55), Inches(y), Inches(8.9), Inches(1.15))
        card.fill.solid()
        card.fill.fore_color.rgb = LIGHT_BG
        card.line.color.rgb = PRIMARY
        tf = card.text_frame
        tf.margin_left = Inches(0.15)
        p0 = tf.paragraphs[0]
        p0.text = f"{name}  ·  {term}  ·  {price}"
        p0.font.size = Pt(17)
        p0.font.bold = True
        p0.font.color.rgb = PRIMARY
        p1 = tf.add_paragraph()
        p1.text = f"包含：{scope}"
        p1.font.size = Pt(13)
        p1.font.color.rgb = DARK
        y += 1.28


def slide_process(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_title_bar(slide, "签约后 3 步上手", "降低客户决策顾虑")
    steps = [
        ("第 1 步", "1～2 天", "安装系统 + 登录配置 + 路径/规则初始化"),
        ("第 2 步", "第 3～5 天", "试跑 10～20 个 SKU，确认发品与监控正常"),
        ("第 3 步", "第 2 周起", "开启定时发品 + 周度数据分析例会"),
    ]
    x = 0.55
    for title, days, desc in steps:
        card = slide.shapes.add_shape(1, Inches(x), Inches(2.0), Inches(2.85), Inches(3.8))
        card.fill.solid()
        card.fill.fore_color.rgb = LIGHT_BG
        card.line.color.rgb = ACCENT
        tf = card.text_frame
        tf.margin_top = Inches(0.25)
        p0 = tf.paragraphs[0]
        p0.text = title
        p0.font.size = Pt(22)
        p0.font.bold = True
        p0.font.color.rgb = ACCENT
        p0.alignment = PP_ALIGN.CENTER
        p1 = tf.add_paragraph()
        p1.text = days
        p1.font.size = Pt(16)
        p1.font.color.rgb = PRIMARY
        p1.alignment = PP_ALIGN.CENTER
        p1.space_before = Pt(12)
        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.size = Pt(13)
        p2.font.color.rgb = DARK
        p2.alignment = PP_ALIGN.CENTER
        p2.space_before = Pt(16)
        x += 3.1
    add_highlight_box(slide, "销售承诺模板：____ 个工作日内完成首批发品成功。", top=6.0)


def slide_case(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_title_bar(slide, "客户案例（待补充）", "有案例后成交率显著提升")
    add_bullets(
        slide,
        [
            "【案例 1】行业：________  |  规模：____ SKU",
            "  • 使用前痛点：________________________________",
            "  • 使用后效果：发品效率 ____  |  询盘变化 ____",
            "  • 客户原话：「________________________________」",
            "",
            "【案例 2】行业：________  |  规模：____ SKU",
            "  • 使用前痛点：________________________________",
            "  • 使用后效果：________________________________",
            "",
            "【案例 3】可贴聊天截图 / 数据对比图（建议放本页右侧）",
        ],
        top=1.7,
        font_size=16,
    )
    hint = slide.shapes.add_shape(1, Inches(6.0), Inches(2.2), Inches(3.4), Inches(4.2))
    hint.fill.solid()
    hint.fill.fore_color.rgb = RGBColor(0xE8, 0xEC, 0xF0)
    hint.line.color.rgb = MUTED
    hint.text_frame.text = "截图\n占位区"
    hint.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    hint.text_frame.paragraphs[0].font.size = Pt(18)
    hint.text_frame.paragraphs[0].font.color.rgb = MUTED


def slide_demo_flow(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_title_bar(slide, "现场演示建议流程（15 分钟）", "销售拜访 / 线上演示脚本")
    add_bullets(
        slide,
        [
            "① 控制台（2 分钟）— 展示异动、P4P、推荐关注，建立「有决策依据」印象",
            "② 自动发品（5 分钟）— 展示七步流程、定时、实时日志，建立「能落地」印象",
            "③ 新发监控（3 分钟）— 展示发品后自动入库 + 周曝光筛选",
            "④ 数据下载+分析（3 分钟）— 一键任务 + 综合分析结果",
            "⑤ 收尾（2 分钟）— 回到报价页 + 试用/签约政策",
        ],
        top=1.75,
        font_size=17,
    )
    add_highlight_box(slide, "演示前务必：登录有效、准备 1 份真实 Excel、确认网络与浏览器正常。")


def slide_objections(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_title_bar(slide, "常见拒绝与应对", "销售话术备忘")
    add_bullets(
        slide,
        [
            "「太贵了」→ 算人力账：1 个运营月薪 vs 系统年费；再问「漏发/错价一次的成本」",
            "「我们自己有脚本」→ 问维护谁、是否含分析看板、新发是否自动跟踪",
            "「先观望」→ 提供试用 + 陪跑 10 条 SKU，用结果说话",
            "「怕封号」→ 强调合规使用、间隔可配置、模拟真人节奏，不承诺排名",
            "「老板不在」→ 留一页「老板版」：三大价值 + 对比表 + 报价，约二次会议",
        ],
        top=1.75,
        font_size=16,
    )


def slide_cta(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, PRIMARY)
    t1 = slide.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(8.4), Inches(0.9))
    t1.text_frame.text = "下一步行动"
    t1.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    t1.text_frame.paragraphs[0].font.size = Pt(36)
    t1.text_frame.paragraphs[0].font.bold = True
    t1.text_frame.paragraphs[0].font.color.rgb = WHITE

    actions = [
        "① 预约 15 分钟在线演示",
        "② 申请试用账号（____ 天）",
        "③ 确定合作方案与付款方式",
        "④ 安排安装陪跑，首周见效果",
    ]
    box = slide.shapes.add_textbox(Inches(1.2), Inches(3.2), Inches(7.6), Inches(2.2))
    tf = box.text_frame
    for i, a in enumerate(actions):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = a
        p.font.size = Pt(20)
        p.font.color.rgb = RGBColor(0xE0, 0xEE, 0xFA)
        p.space_after = Pt(12)
        p.alignment = PP_ALIGN.CENTER

    contact = slide.shapes.add_textbox(Inches(0.8), Inches(5.5), Inches(8.4), Inches(1.2))
    ctf = contact.text_frame
    ctf.text = (
        "【请填写】\n"
        "销售顾问：________    手机/微信：________\n"
        "公司地址：________    官网/演示链接：________"
    )
    for p in ctf.paragraphs:
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(15)
        p.font.color.rgb = RGBColor(0xB0, 0xCC, 0xE4)


def slide_end(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, PRIMARY)
    t1 = slide.shapes.add_textbox(Inches(1), Inches(2.8), Inches(8), Inches(1.2))
    tf = t1.text_frame
    tf.text = "让团队把时间花在询盘和订单上"
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].font.size = Pt(32)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = WHITE

    t2 = slide.shapes.add_textbox(Inches(1), Inches(4.2), Inches(8), Inches(0.8))
    tf2 = t2.text_frame
    tf2.text = "感谢您的信任 · 期待合作"
    tf2.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf2.paragraphs[0].font.size = Pt(20)
    tf2.paragraphs[0].font.color.rgb = RGBColor(0xC8, 0xDC, 0xEC)


def build():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    slide_cover(prs)
    slide_hook(prs)
    slide_solution(prs)
    slide_before_after(prs)
    slide_value_pillars(prs)
    slide_product_publish(prs)
    slide_product_data(prs)
    slide_product_assets(prs)
    slide_dashboard(prs)
    slide_scenarios(prs)
    slide_compare(prs)
    slide_trust(prs)
    slide_delivery(prs)
    slide_pricing(prs)
    slide_process(prs)
    slide_case(prs)
    slide_demo_flow(prs)
    slide_objections(prs)
    slide_cta(prs)
    slide_end(prs)

    OUTPUT_CN.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUTPUT_CN))
    prs.save(str(OUTPUT_EN))
    print(f"已生成: {OUTPUT_CN}")
    print(f"已生成: {OUTPUT_EN}")
    print(f"共 {len(prs.slides)} 页")


if __name__ == "__main__":
    build()
